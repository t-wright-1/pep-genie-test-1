    
from asyncore import file_dispatcher
import cv2
import numpy as np
import ast
import os
import zipfile
from PIL import Image, ImageStat, ImageDraw, ImageOps
import matplotlib.pyplot as plt

from django.conf import settings
media_root = settings.MEDIA_ROOT

from pptx import Presentation
from pptx.util import Cm


#Two functions to get four-point transformed image from blank photo and coordinates string
def order_points(pts):
    rect = np.zeros((4, 2), dtype = "float32")

    s = pts.sum(axis = 1)
    rect[0] = pts[np.argmin(s)]
    rect[2] = pts[np.argmax(s)]

    diff = np.diff(pts, axis = 1)
    rect[1] = pts[np.argmin(diff)]
    rect[3] = pts[np.argmax(diff)]
    print(rect)
    return rect

def four_point_transform(image_path, in_str):
    
    print(os.path.exists(image_path))
    print(image_path)

    #Convert string of fractions to list of coords
    image = cv2.imread(image_path)
    dims = image.shape

    my_str = ast.literal_eval(in_str)
    my_row = []
    for i in my_str:
        my_row.append(i[0]*dims[1])
        my_row.append(i[1]*dims[0])
        my_str[my_str.index(i)] = my_row
        my_row = []
    my_str = np.array(my_str)
    pts = order_points(my_str) #list of ordered coords

    #Transorm image using list of coords (pts) and image
    (tl, tr, br, bl) = pts

    widthA = np.sqrt(((br[0] - bl[0]) ** 2) + ((br[1] - bl[1]) ** 2))
    widthB = np.sqrt(((tr[0] - tl[0]) ** 2) + ((tr[1] - tl[1]) ** 2))
    maxWidth = max(int(widthA), int(widthB))

    heightA = np.sqrt(((tr[0] - br[0]) ** 2) + ((tr[1] - br[1]) ** 2))
    heightB = np.sqrt(((tl[0] - bl[0]) ** 2) + ((tl[1] - bl[1]) ** 2))
    maxHeight = max(int(heightA), int(heightB))

    dst = np.array([
        [0, 0],
        [maxWidth - 1, 0],
        [maxWidth - 1, maxHeight - 1],
        [0, maxHeight - 1]], dtype = "float32")

    M = cv2.getPerspectiveTransform(pts, dst)
    warped = cv2.warpPerspective(image, M, (maxWidth, maxHeight))     # return the warped image
    warped = cv2.rotate(warped, cv2.ROTATE_90_COUNTERCLOCKWISE)
    warped = cv2.flip(warped, 1)
    return warped


def strip_number(input_string, first_col):
    
    input_string = input_string.replace(' ','')
    input_list = input_string.split(',')

    letter_ref = ['7','A', 'B', 'C', 'D', 'E', 'F', 'G', 'H', 'I', 'J', 'K', 'L', 'M', 'N', 'O', 'P', 'Q', 'R', 'S', 'T', 'U', 'V', 'W', 'X', 'Y', 'Z']

    new_strip = []
    strip_list=[]
    for strip in input_list:
 
        #convert letter to row number
        if letter_ref.index(strip[0]) < 10:
            new_strip.append('0')
            new_strip.append(letter_ref.index(strip[0]))
        else:
            new_strip.append(letter_ref.index(strip[0]))
        #add column number
        new_strip.append(strip[1:3])
        new_strip.append('-')

        #convert letter 2 to row number
        if letter_ref.index(strip[4]) < 10:
            new_strip.append('0')
            new_strip.append(letter_ref.index(strip[4]))
        else:
            new_strip.append(letter_ref.index(strip[4]))
        #add column number
        new_strip.append(strip[5:7])

        #add output to list
        strip_list.append(''.join(map(str, new_strip)))
        new_strip = []


    new_item = 0
    new_row = []
    result=[]
    for i in strip_list:
        #add row number
        new_item += int(i[0:2])
        #add col number
        new_item += (int(i[2:4])-first_col)*20
        new_row.append(new_item)
        new_item = 0

        new_item += int(i[5:7])
        new_item += (int(i[7:9])-first_col)*20
        new_row.append(new_item)
        new_item = 0
        result.append(new_row)
        new_row = []
    #for i in result:
        #print(i[0], i[1])

    return result


def strip_maker(image_path, rows, cols, strips):
    #variables
    my_image = Image.open(image_path)
    old_width, old_height = my_image.size

    #resize cropped image
    if old_height <= 400:
        new_width = int( (400 / old_height) * old_width )
        my_image = my_image.resize((new_width, 400))

    #resized image variables
    width, height = my_image.size
    slice_height = height / rows
    slice_width = width / cols
    upper = 0
    left = 0

    #cut cropped image into squares/slices
    count = 1
    slice_list = []
    for col in range(cols):
        right = (col+1)* slice_width
        for row in range(rows):
            if count == rows:
                lower = height
            else:
                lower = int(count * slice_height)
            bbox = (left, upper, right, lower)
            working_slice = my_image.crop(bbox)
            upper += int(slice_height)
            slice_list.append(working_slice)
            my_width, my_height = working_slice.size
            #print(count, my_height, lower, upper)
            count += 1
        left += slice_width
        upper = 0
        count = 1

    # convert input string to number tuple
    output_strips = []
    for number_pair in strips:
        new_strip = []
        new_strip.append(number_pair[0]-1)
        new_strip.append(number_pair[1])
        image_number = new_strip[1]-new_strip[0]

        #Make strip canvas and join squares/slices
        width, height = slice_list[0].size
        new_image = Image.new('RGB', (width, height * image_number))
        counter = 0
        for image in slice_list[new_strip[0]:new_strip[1]]:
            new_image.paste(image,(0,counter*height))
            counter +=1
        output_strips.append(new_image) 
        width, height = output_strips[0].size 

    return output_strips

def make_ppt(test_list, mock_list):
    prs = Presentation()
    
    left = Cm(1)
    top = Cm(1)
    path = os.path.join(media_root,'mock_crop.jpg')

    for i in range(len(test_list)):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        img=slide.shapes.add_picture(path,left,top)


    prs.save(os.path.join(media_root,'presentation1.pptx'))


def slice_to_list(my_image, rows, cols):
    #variables
    #my_image = Image.open(image_path)
    old_width, old_height = my_image.size #was issue opening image from model, but 
                                        # now I am supplying the PIL file

    #resize cropped image
    if old_height <= 400:
        new_width = int( (400 / old_height) * old_width )
        my_image = my_image.resize((new_width, 400))

    #resized image variables
    width, height = my_image.size
    slice_height = height / rows
    slice_width = width / cols
    upper = 0
    left = 0

    #cut cropped image into squares/slices
    count = 1
    slice_list = []
    for col in range(cols):
        right = (col+1)* slice_width
        for row in range(rows):
            if count == rows:
                lower = height
            else:
                lower = int(count * slice_height)
            bbox = (left, upper, right, lower)
            working_slice = my_image.crop(bbox)
            upper += int(slice_height)
            slice_list.append(working_slice)
            my_width, my_height = working_slice.size
            #print(count, my_height, lower, upper)
            count += 1
        left += slice_width
        upper = 0
        count = 1
    return slice_list








#Crop image to circle and measure

#mask variables
#width = 200
def make_mask(width):
    frac = width / 15

    #create mask image
    bbox = [frac, frac, width-frac, width-frac]
    mask = Image.new('L', (width, width),0)
    draw = ImageDraw.Draw(mask)
    draw.pieslice(bbox,0,360,fill=255)
    return mask

def circle_crop(im, mask):
    #crop image to mask
    cropped = ImageOps.fit(im, mask.size, centering=(0.5, 0.5))
    #add alpha layer to image
    cropped.putalpha(mask)
    #output is RGBA
    return cropped

def measure_darkness(im):
    stat = ImageStat.Stat(im)
    return stat.mean[0]



   




def stripper2(image_path, graph_type_list, rows, cols, strips):
    #variables
    my_image = Image.open(image_path)
    old_width, old_height = my_image.size

    #resize cropped image
    new_width = int( (1000 / old_height) * old_width )
    my_image = my_image.resize((new_width, 1000))

    #resized image variables
    width, height = my_image.size
    slice_height = height / rows
    slice_width = width / cols
    upper = 0
    left = 0

    #cut cropped image into squares/slices
    count = 1
    slice_list = []
    for col in range(cols):
        right = (col+1)* slice_width
        for row in range(rows):
            if count == rows:
                lower = height
            else:
                lower = int(count * slice_height)
            bbox = (left, upper, right, lower)
            working_slice = my_image.crop(bbox)
            upper += int(slice_height)
            slice_list.append(working_slice)
            my_width, my_height = working_slice.size
            #print(count, my_height, lower, upper)
            count += 1
        left += slice_width
        upper = 0
        count = 1

    # convert input string to number tuple
    output_strips = []
    for number_pair in strips:
        new_strip = []
        new_strip.append(number_pair[0]-1)
        new_strip.append(number_pair[1])
        image_number = new_strip[1]-new_strip[0]

        #Make strip canvas and join squares/slices
        width, height = slice_list[0].size

        new_image = Image.new('RGB', (width, height * image_number))

        counter = 0
        index_counter = new_strip[0]+1
        for image in slice_list[new_strip[0]:new_strip[1]]:
            new_image.paste(image,(0,counter*height))
            if ((index_counter-1)/20).is_integer() and counter != 0:
                xy = (0,counter*height,width,counter*height)
                draw = ImageDraw.Draw(new_image)
                draw.line(xy, fill='black', width=3)
            index_counter += 1
            counter +=1
        if graph_type_list[strips.index(number_pair)] == 'v':
            new_image = new_image.transpose(Image.ROTATE_90)
            new_image = new_image.resize((height*image_number,width)) 
        output_strips.append(new_image) 
        width, height = output_strips[0].size 


    return output_strips









### MAKE GRAPHS ###---------------------------------------


def h_bar(my_x,my_heights,my_dir,my_filename,ylim=False):
    plt.clf()
    plt.barh(my_x,my_heights)

    #invert y
    ax = plt.gca() #get current axis
    ax.invert_yaxis()

    #hide y ticks
    ax.yaxis.set_ticklabels([])
    plt.yticks([], [])
    

    #add threshold
    if my_x[0] != 0:
        plt.axvline(x=1,color='blue',lw=1, ls='--')

    #hide frame/border
    ax.spines['bottom'].set_visible(False)
    ax.spines['right'].set_visible(False)

    #ax.set_xlabel('Relative Luminescence',fontsize=14)

    #x axis on top
    ax.xaxis.tick_top()


    plt.margins(y=0)

    #set axis if disc array
    if ylim != False:
        plt.xlim(0,ylim)

    #plt.show()
    fig = plt.gcf()
    spot_no = len(my_x)
    width = (7.22/20)*spot_no
    fig.set_size_inches(5, width) #8.37/20 now 7.2/20 per seq, width 6

    my_path = os.path.join(my_dir,my_filename)
    fig.savefig(my_path, bbox_inches='tight',dpi=100)






def v_bar(my_x,my_heights,my_dir,my_filename,ylim=False):
    plt.clf()
    plt.bar(my_x,my_heights)

    ax = plt.gca() #get current axis    

    #add threshold
    if my_x[0] != 0:
        plt.axhline(y=1,color='blue',lw=1, ls='--')

    #hide frame/border
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    #hide y ticks
    ax.xaxis.set_ticklabels([])
    plt.xticks([], [])

    #ax.set_xlabel('Relative Luminescence',fontsize=14)

    plt.margins(x=0)
    #plt.show()

    #set axis if disc array
    if ylim != False:
        plt.ylim(0,ylim)

    fig = plt.gcf()
    spot_no = len(my_x)
    width = (7.22/20)*spot_no
    fig.set_size_inches(width,4) #8.37/20 now 7.2/20 per seq, width 6
    my_path = os.path.join(my_dir,my_filename)
    fig.savefig(my_path, bbox_inches='tight',dpi=100)


def graph_to_ppt(path_list,graph_type_list,ppt_dir='blank'):
    if ppt_dir == 'blank':
        prs = Presentation()
        prs.slide_width = Cm(33.9)
        prs.slide_height = Cm(19.1)
        for i in range(len(path_list)):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
    else:
        prs = Presentation(ppt_dir)
    
    h_left = Cm(8.7)
    h_top = Cm(2.2)
    v_left = Cm(8.4)
    v_top = Cm(6.5)
    #width = Cm(0.7)
    
    slides_list = [slide.slide_id for slide in prs.slides]

    c=0
    for i in path_list:
        #slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide = prs.slides.get(slides_list[c])
        if graph_type_list[c] == "h":
            img = slide.shapes.add_picture(i,h_left,h_top)
        if graph_type_list[c] == "v":
            img = slide.shapes.add_picture(i,v_left,v_top)
        c+=1
    prs.save('graph_ppt')




def file_compress(inp_file_names, filenames, out_zip_file):
    compression = zipfile.ZIP_DEFLATED
    zf = zipfile.ZipFile(out_zip_file, mode="w")
    try:
        for i in range(len(inp_file_names)):
    
            zf.write(inp_file_names[i], filenames[i], compress_type=compression)

    except FileNotFoundError as e:
        print('compression error')
    finally:
        zf.close()

