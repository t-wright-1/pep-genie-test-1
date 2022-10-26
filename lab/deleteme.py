import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Cm

#invert values so darkness is higher value, and normalise relative to first spot
def norm_list(my_list):
    my_min = min(my_list)
    my_max = max(my_list)
    my_range = my_max-my_min
    def norm_invert(i, my_min, my_range):
        i = ( (1- ((i-my_min)/my_range)  ) * my_range ) +my_min
        return i
    new_list = [norm_invert(i,my_min,my_range) for i in my_list]
    new_list_2 = [i/new_list[0] for i in new_list]
    return new_list_2







def h_bar(my_items,my_x):
    plt.clf()
    plt.barh(my_items,my_x)

    #invert y
    ax = plt.gca() #get current axis
    ax.invert_yaxis()

    #hide y ticks
    ax.yaxis.set_ticklabels([])
    plt.yticks([], [])
    

    #add threshold
    if my_x[0] != 0:
        plt.axvline(x=1,color='blue',lw=1, ls='--')

    #hide frame/border
    ax.spines['bottom'].set_visible(False)
    ax.spines['right'].set_visible(False)

    #ax.set_xlabel('Relative Luminescence',fontsize=14)

    #x axis on top
    ax.xaxis.tick_top()


    plt.margins(y=0)

    #plt.show()
    fig = plt.gcf()
    spot_no = len(my_x)
    width = (7.22/20)*spot_no
    fig.set_size_inches(5, width) #8.37/20 now 7.2/20 per seq, width 6
    fig.savefig('graph_test_h.png', bbox_inches='tight',dpi=100)



def v_bar(my_items,my_x):
    plt.clf()
    plt.bar(my_items,my_x)

    ax = plt.gca() #get current axis    

    #add threshold
    if my_x[0] != 0:
        plt.axhline(y=1,color='blue',lw=1, ls='--')

    #hide frame/border
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    #hide y ticks
    ax.xaxis.set_ticklabels([])
    plt.xticks([], [])

    #ax.set_xlabel('Relative Luminescence',fontsize=14)

    plt.margins(x=0)
    #plt.show()

    fig = plt.gcf()
    spot_no = len(my_x)
    width = (7.22/20)*spot_no
    fig.set_size_inches(width,4) #8.37/20 now 7.2/20 per seq, width 6
    fig.savefig('graph_test_v.png', bbox_inches='tight',dpi=100)


def graph_to_ppt(path_list,graph_type_list,ppt_dir='blank'):
    if ppt_dir == 'blank':
        prs = Presentation()
        prs.slide_width = Cm(33.9)
        prs.slide_height = Cm(19.1)
        for i in range(len(path_list)):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
    else:
        prs = Presentation(ppt_dir)
    
    h_left = Cm(8.7)
    h_top = Cm(2.2)
    v_left = Cm(8.4)
    v_top = Cm(6.5)
    #width = Cm(0.7)
    
    slides_list = [slide.slide_id for slide in prs.slides]

    c=0
    for i in path_list:
        #slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide = prs.slides.get(slides_list[c])
        if graph_type_list[c] == "h":
            img = slide.shapes.add_picture(i,h_left,h_top)
        if graph_type_list[c] == "v":
            img = slide.shapes.add_picture(i,v_left,v_top)
        c+=1
    prs.save('graph_ppt')

my_x = [2,3,4,5,6,7,8,7,6,5,4,2,3,4,5,6,7,8,7,6]
my_y = range(len(my_x))

v_bar(my_y,norm_list(my_x))
h_bar(my_y,norm_list(my_x))

path_list = ['graph_test_v.png']
graph_type_list = ['v','v']

graph_to_ppt(path_list,graph_type_list,'hoz.pptx')
#graph_to_ppt(path_list,'the_big_pres.pptx')

#Next
#Add relative luminescence label
#remove weird title slide
#Add baillie logo with last slide - instructions for text
#Make django infrastructure to import slide
#Correct normalisation of data.