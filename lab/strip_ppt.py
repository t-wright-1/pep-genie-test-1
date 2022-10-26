from dis import dis
from PIL import Image, ImageDraw
from pptx import Presentation 
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

from django.conf import settings
media_root = settings.MEDIA_ROOT

def stripper(image_path, graph_type_list, rows, cols, strips):
    #variables
    my_image = Image.open(image_path)
    old_width, old_height = my_image.size

    #resize cropped image
    new_width = int( (1000 / old_height) * old_width )
    my_image = my_image.resize((new_width, 1000))

    #resized image variables
    width, height = my_image.size
    slice_height = height / rows
    slice_width = width / cols
    upper = 0
    left = 0

    #cut cropped image into squares/slices
    count = 1
    slice_list = []
    for col in range(cols):
        right = (col+1)* slice_width
        for row in range(rows):
            if count == rows:
                lower = height
            else:
                lower = int(count * slice_height)
            bbox = (left, upper, right, lower)
            working_slice = my_image.crop(bbox)
            upper += int(slice_height)
            slice_list.append(working_slice)
            my_width, my_height = working_slice.size
            #print(count, my_height, lower, upper)
            count += 1
        left += slice_width
        upper = 0
        count = 1

    # convert input string to number tuple
    output_strips = []
    for number_pair in strips:
        new_strip = []
        new_strip.append(number_pair[0]-1)
        new_strip.append(number_pair[1])
        image_number = new_strip[1]-new_strip[0]

        #Make strip canvas and join squares/slices
        width, height = slice_list[0].size

        new_image = Image.new('RGB', (width, height * image_number))

        counter = 0
        index_counter = new_strip[0]+1
        for image in slice_list[new_strip[0]:new_strip[1]]:
            new_image.paste(image,(0,counter*height))
            if ((index_counter-1)/20).is_integer() and counter != 0:
                xy = (0,counter*height,width,counter*height)
                draw = ImageDraw.Draw(new_image)
                draw.line(xy, fill='black', width=3)
            index_counter += 1
            counter +=1
        if graph_type_list[strips.index(number_pair)] == 'v':
            new_image = new_image.transpose(Image.ROTATE_90)
            new_image = new_image.resize((height*image_number,width)) 
            print(new_image.size, graph_type_list)
        output_strips.append(new_image) 
        width, height = output_strips[0].size 

    return output_strips

def make_pptx(test_list, graph_type_list,strip_request,user_dir,user_id):
    prs = Presentation()
    
    prs.slide_width = Cm(33.9)
    prs.slide_height = Cm(19.1)
    



    for i in range(len(test_list)):
        #get number of slices for height calculation
        new_strip = []
        new_strip.append(strip_request[i][0]-1)
        new_strip.append(strip_request[i][1])
        image_number = new_strip[1]-new_strip[0]
        long = Cm(((20.3*0.7)/20)*image_number)
        short = Cm(0.7)
        #account for vertical vs horizontal
        if graph_type_list[i] == 'h': #functional
            width = short
            height = long
            left = Cm(8) #one cm before
            left_mock = Cm(7)
            top = Cm(3)
            top_mock = Cm(3)
            
            text_left = Cm(6.75)
            text_top = Cm(2)
            text_width = Cm(1)
            text_height = Cm(1)
            text_align = PP_ALIGN.LEFT
            text_rotation = -90.0
            text_left_factor = Cm(1.1)
            text_top_factor = Cm(0)

        elif graph_type_list[i] == 'v': 
            width = long
            height = short
            left = Cm(9.5) #one cm before
            left_mock = Cm(9.5)
            top = Cm(15) #
            top_mock = Cm(16) #
 
            text_left = Cm(2.5)
            text_top = Cm(14.9)
            text_width = Cm(7)
            text_height = Cm(1)
            text_align = PP_ALIGN.RIGHT
            text_rotation = 0
            text_left_factor = Cm(0)
            text_top_factor = Cm(0.9)

        slide = prs.slides.add_slide(prs.slide_layouts[5])

        my_file_name = 'strip_' + str(i) + '.png'
        my_path = os.path.join(user_dir, my_file_name)
        print(width,height)
        img=slide.shapes.add_picture(my_path,left,top,width,height)

        my_file_name_mock = 'mock_strip_' + str(i) + '.png'
        my_path_mock = os.path.join(user_dir, my_file_name_mock)
        img = slide.shapes.add_picture(my_path_mock,left_mock,top_mock,width,height)

        #add labels for strips
        txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        tf = txBox.text_frame
        tf.text = "Control"
        txBox.rotation = text_rotation
        txBox.text_frame.paragraphs[0].alignment = text_align
        
        txBox2 = slide.shapes.add_textbox(text_left+text_left_factor, text_top+text_top_factor, text_width, text_height)
        tf2 = txBox2.text_frame
        tf2.text = "Overlay"
        txBox2.rotation = text_rotation
        txBox2.text_frame.paragraphs[0].alignment = text_align

    my_filename = 'presentation_'+user_id+'.pptx'
    prs.save(os.path.join(user_dir,my_filename))

#second function customised for sd
