from dis import dis
from PIL import Image
from pptx import Presentation
from pptx.util import Cm
import os

from django.conf import settings
media_root = settings.MEDIA_ROOT

def stripper(image_path, rows, cols, strips):
    #variables
    my_image = Image.open(image_path)
    old_width, old_height = my_image.size

    #resize cropped image
    new_width = int( (1000 / old_height) * old_width )
    my_image = my_image.resize((new_width, 1000))

    #resized image variables
    width, height = my_image.size
    slice_height = height / rows
    slice_width = width / cols
    upper = 0
    left = 0

    #cut cropped image into squares/slices
    count = 1
    slice_list = []
    for col in range(cols):
        right = (col+1)* slice_width
        for row in range(rows):
            if count == rows:
                lower = height
            else:
                lower = int(count * slice_height)
            bbox = (left, upper, right, lower)
            working_slice = my_image.crop(bbox)
            upper += int(slice_height)
            slice_list.append(working_slice)
            my_width, my_height = working_slice.size
            #print(count, my_height, lower, upper)
            count += 1
        left += slice_width
        upper = 0
        count = 1

    # convert input string to number tuple
    output_strips = []
    for number_pair in strips:
        new_strip = []
        new_strip.append(number_pair[0]-1)
        new_strip.append(number_pair[1])
        image_number = new_strip[1]-new_strip[0]

        #Disconnect detector: if images are spliced, border needs to be added
        disconnect = False
        for i in range(new_strip[0],(new_strip[1]-1)):
            if (i/20).is_integer():
                disconnect = True
        #Make strip canvas and join squares/slices
        width, height = slice_list[0].size
        #set canvas height and add space for border if disconnect
        if disconnect == True:
            new_image = Image.new('RGB', (width, (height * image_number)+(int(height/6))))
        else: 
            new_image = Image.new('RGB', (width, height * image_number))

        counter = 0
        border_shift = 0

        for image in slice_list[new_strip[0]:new_strip[1]]:
            #if border needed
            if disconnect == True:
                if (slice_list.index(image)/20).is_integer():
                    border_shift = int(height/6)
            new_image.paste(image,(0,counter*height+border_shift))
            counter +=1
        output_strips.append(new_image) 
        width, height = output_strips[0].size 

    return output_strips

def make_pptx(test_list,strip_request,user_dir,user_id):
    prs = Presentation()
    
    prs.slide_width = Cm(33.9)
    prs.slide_height = Cm(19.1)
    left = Cm(2)
    left_mock = Cm(1)
    top = Cm(1)
    width = Cm(0.7)

    for i in range(len(test_list)):
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        title_placeholder = slide.shapes.title
        title_placeholder.text = str(strip_request[i])

        my_file_name = 'strip_' + str(i) + '.png'
        my_path = os.path.join(user_dir, my_file_name)
        img=slide.shapes.add_picture(my_path,left,top,width)

        my_file_name_mock = 'mock_strip_' + str(i) + '.png'
        my_path_mock = os.path.join(user_dir, my_file_name_mock)
        img=slide.shapes.add_picture(my_path_mock,left_mock,top,width)
    my_filename = 'presentation_'+user_id+'.pptx'
    prs.save(os.path.join(user_dir,my_filename))
